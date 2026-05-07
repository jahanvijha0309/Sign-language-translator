import os
from pptx import Presentation
from pptx.util import Inches, Pt

# Output in same folder as this script - works on any OS
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(BASE_DIR, "Sign_Language_Presentation.pptx")

slides = [
    {
        "title": "Sign Language to Text Translator",
        "content": ["End-to-end demo: Flask + TensorFlow + OpenCV", "Author: (Your Name) — Date"],
        "notes": "Opening slide. Introduce yourself and the project goal."
    },
    {
        "title": "Motivation",
        "content": ["Improve accessibility for deaf users", "Reduce reliance on human interpreters for simple messages"],
        "notes": "Explain the problem and target users."
    },
    {
        "title": "Goal & Requirements",
        "content": ["Real-time webcam translation", "Upload video/image support", "Web-accessible and deployable (Docker)"],
        "notes": "State the project goals and constraints."
    },
    {
        "title": "System Overview",
        "content": ["Browser UI <-> Flask backend <-> Keras model", "Optional MediaPipe hand detection for cropping"],
        "notes": "Mention routes: /predict_image, /predict_video, /predict_live"
    },
    {
        "title": "Dataset & Labels",
        "content": ["Folder-per-class structure (A..Z, space, delete, nothing)", "Augmentation via ImageDataGenerator"],
        "notes": "Show sample images when presenting."
    },
    {
        "title": "Model & Training",
        "content": ["Transfer learning: MobileNetV2 base", "Final dense layer size = number of labels", "Model saved as modelnet_model.h5"],
        "notes": "Explain choice of transfer learning and IMG_SIZE=224."
    },
    {
        "title": "Demo / UX",
        "content": ["Upload video/image or use Live Webcam", "Realtime overlay of predicted label and confidence"],
        "notes": "Walk through the UI and show live demo steps."
    },
    {
        "title": "Results & Limitations",
        "content": ["Good for isolated letters; sentence decoding is future work", "Limitations: lighting, occlusion, multi-hand scenarios"],
        "notes": "Be honest about limitations and trade-offs."
    },
    {
        "title": "Future Work",
        "content": ["Sequence modeling for sentences", "Robust MediaPipe integration and more data", "Cloud deployment / mobile app"],
        "notes": "Suggest next steps and invite collaboration."
    },
    {
        "title": "Thank You / Q&A",
        "content": ["Links: README, GitHub repo, demo video", "Contact info"],
        "notes": "End with call-to-action and invite questions."
    }
]


def add_bullet_slide(prs, title, bullets, notes_text=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def build_presentation(output_path=OUTPUT_PATH):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sign Language to Text Translator"
    slide.placeholders[1].text = "Flask + TensorFlow + OpenCV — demo"

    for s in slides:
        add_bullet_slide(prs, s["title"], s["content"], s.get("notes"))

    prs.save(output_path)
    print(f"[OK] Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
